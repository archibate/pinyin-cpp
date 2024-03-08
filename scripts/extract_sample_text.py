import os
import random
import tiktoken
import glob
import json
import pptx
from multiprocessing import Pool, cpu_count
import pypinyin

files = []
this_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)))
with open(os.path.join(this_dir, "sample_text_files.txt")) as fin:
    for line in fin:
        line = line.strip()
        if line:
            files += glob.glob(line)

def process_file(file):
    strings = []
    basename, extname = os.path.splitext(file)
    strings.append(basename)
    if extname in '.md':
        with open(file) as md:
            for line in md:
                line = line.strip()
                if line:
                    strings.append(line)
    elif extname in '.ass':
        with open(file) as md:
            for line in md:
                line = line.strip()
                if line.startswith('Dialogue:'):
                    strings.append(line.split(',,')[-1])
    elif extname in '.srt':
        with open(file) as md:
            for line in md:
                line = line.strip()
                if line and not line[0].isdigit():
                    strings.append(line)
    elif extname == '.pptx':
        p = pptx.Presentation(file)
        for slide in p.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    strings.append(shape.text)
    return '\n'.join(s for s in map(str.strip, strings) if s)

def chinese_rate(text):
    score = -0.1
    for c in text:
        if 0x4e00 <= ord(c) <= 0x9fff:
            score += 1
        elif c.isalpha():
            score -= 0.1
    return score

def process_line(line):
    line = line.replace('\r', ' ').replace('\n', ' ').replace(' ', '').strip()
    if chinese_rate(line) < 0:
        return ''
    return line

with Pool(cpu_count()) as pool:
    result = '\n'.join(pool.map(process_file, files)).splitlines()
    random.shuffle(result)
    result = result[:round(len(result) * 1)]
    result = '\n'.join(x for x in pool.map(process_line, result) if x)

with open("/tmp/extract.txt", "w") as fout:
    fout.write(result)
